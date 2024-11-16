from pptx import Presentation
from pptx.util import Inches, Pt
# from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def set_font(para, font_name, font_size, font_color=RGBColor(0, 0, 0)):
    run = para.add_run()
    font = run.font
    font.name = font_name
    font.size = font_size
    # para.alignment = PP_ALIGN.CENTER
    run.font.color.rgb = font_color
    return para, run

def add_title_slide(prs, title):
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title_placeholder = title_slide.shapes.title
    subtitle_placeholder = title_slide.placeholders[1]

    title_placeholder.text = title
    subtitle_placeholder.text = "Created with Python and python-pptx"

def add_content_slide(prs, title, content):
    content_slide_layout = prs.slide_layouts[1]
    content_slide = prs.slides.add_slide(content_slide_layout)

    title_placeholder = content_slide.shapes.title
    content_placeholder = content_slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

def main():
    presentation = Presentation()

    # Add a title slide
    add_title_slide(presentation, "Awesome Python Presentation")

    # Add content slides
    slide_contents = [
        ("Slide 1", "This is the content for Slide 1."),
        ("Slide 2", "This is the content for Slide 2."),
        # Add more slides as needed
    ]

    for title, content in slide_contents:
        add_content_slide(presentation, title, content)

    # Save the presentation
    presentation.save("awesome_presentation.pptx")

if __name__ == "__main__":
    main()
